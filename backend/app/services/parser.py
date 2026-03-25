import fitz                        # PyMuPDF
from pptx import Presentation
from docx import Document
from fastapi import HTTPException
import os

def parse_pdf(file_bytes: bytes) -> str:
    """Extract all text from a PDF."""
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text_parts = []
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                text_parts.append(f"[Page {page_num}]\n{text}")
        doc.close()
        return "\n\n".join(text_parts)
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Failed to parse PDF: {str(e)}")


def parse_docx(file_bytes: bytes) -> str:
    """Extract all text from a DOCX file."""
    try:
        import io
        doc = Document(io.BytesIO(file_bytes))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Failed to parse DOCX: {str(e)}")


def parse_pptx(file_bytes: bytes) -> str:
    """Extract all text from a PPTX file, slide by slide."""
    try:
        import io
        prs = Presentation(io.BytesIO(file_bytes))
        slide_texts = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            texts.append(line)
            if texts:
                slide_texts.append(f"[Slide {slide_num}]\n" + "\n".join(texts))
        return "\n\n".join(slide_texts)
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Failed to parse PPTX: {str(e)}")


def parse_document(filename: str, file_bytes: bytes) -> str:
    """Route to the correct parser based on file extension."""
    ext = os.path.splitext(filename)[-1].lower()
    parsers = {
        ".pdf":  parse_pdf,
        ".docx": parse_docx,
        ".pptx": parse_pptx,
    }
    if ext not in parsers:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")

    return parsers[ext](file_bytes)

from typing import List


def parse_pdf_chunks(file_bytes: bytes) -> List[str]:
    """Return PDF text as list of page-wise chunks."""
    full_text = parse_pdf(file_bytes)
    return [chunk.strip() for chunk in full_text.split("\n\n") if chunk.strip()]


def parse_docx_chunks(file_bytes: bytes) -> List[str]:
    """Return DOCX text as list of paragraph chunks."""
    full_text = parse_docx(file_bytes)
    return [chunk.strip() for chunk in full_text.split("\n\n") if chunk.strip()]


def parse_pptx_chunks(file_bytes: bytes) -> List[str]:
    """Return PPTX text as list of slide-wise chunks."""
    full_text = parse_pptx(file_bytes)
    return [chunk.strip() for chunk in full_text.split("\n\n") if chunk.strip()]


def parse_document_chunks(filename: str, file_bytes: bytes) -> List[str]:
    """Return document as List[str] chunks."""
    ext = os.path.splitext(filename)[-1].lower()

    parsers = {
        ".pdf": parse_pdf_chunks,
        ".docx": parse_docx_chunks,
        ".pptx": parse_pptx_chunks,
    }

    if ext not in parsers:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")

    return parsers[ext](file_bytes)